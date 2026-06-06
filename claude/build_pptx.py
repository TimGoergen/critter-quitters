from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

BG       = RGBColor(0xF4, 0xF4, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1A, 0x1A, 0x1A)
MID      = RGBColor(0x55, 0x55, 0x55)
LIGHT    = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT   = RGBColor(0xCC, 0xCC, 0xCC)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_label(slide, text, top):
    tb = slide.shapes.add_textbox(Inches(1.2), top, Inches(10), Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = LIGHT
    run.font.name = "Segoe UI"

def add_heading(slide, text, top, size=36, color=DARK, width=Inches(10.5)):
    tb = slide.shapes.add_textbox(Inches(1.2), top, width, Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Segoe UI"

def add_body(slide, text, top, size=18, color=MID, width=Inches(10.5), italic=False):
    tb = slide.shapes.add_textbox(Inches(1.2), top, width, Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    run.font.italic = italic

def add_bullet_list(slide, items, top, size=18, color=MID):
    tb = slide.shapes.add_textbox(Inches(1.2), top, Inches(10.5), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"

def add_accent_line(slide, top):
    left  = Inches(1.2)
    width = Inches(0.04)
    height = Inches(1.8)
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_time_badge(slide, text):
    tb = slide.shapes.add_textbox(Inches(11.2), Inches(6.8), Inches(1.5), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = LIGHT
    run.font.name = "Segoe UI"

def add_note(slide, text):
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(10.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Note: " + text
    run.font.size = Pt(11)
    run.font.color.rgb = LIGHT
    run.font.italic = True
    run.font.name = "Segoe UI"

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # completely blank layout

# ── Slide 1: Title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_label(s, "FSCA Lightning Talk — 5 minutes", Inches(1.0))
add_heading(s, "Three Goals.\nNone of Which Required Me to Know How.", Inches(1.6), size=40)
add_accent_line(s, Inches(3.6))
items = [
    "Build a mobile game in a platform and language I'd never used.",
    "Build it without writing a single line of code myself.",
    "Get it automatically installed on my phone every time I pushed code.",
]
add_bullet_list(s, items, Inches(3.7), size=17, color=MID)

# ── Slide 2: Beat 1 — The Setup ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_label(s, "Beat 1 — The Setup", Inches(0.9))
add_heading(s, "The Setup", Inches(1.3), size=32)
add_body(s,
    '"I wanted to build a mobile game. I picked Godot — a free, open source game engine '
    "I'd never touched — and GDScript, its scripting language, which I had never written. "
    "I had some Unity experience from years ago, but this was essentially starting from scratch.\"",
    Inches(2.3), size=17, italic=True)
add_note(s, "The point is the unfamiliarity, not the tool. Keep this short.")
add_time_badge(s, "~45 sec")

# ── Slide 3: Beat 2 — The Constraint ────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_label(s, "Beat 2 — The Constraint", Inches(0.9))
add_heading(s, "The Constraint", Inches(1.3), size=32)
add_body(s,
    '"I set a rule for myself: I wasn\'t going to write the code. '
    "I would describe what I wanted, review what came back, and decide whether it was right. "
    'Claude wrote every line. I was the architect."',
    Inches(2.3), size=17, italic=True)
add_note(s, "This is the beat that will get the most reaction. Let it breathe.")
add_time_badge(s, "~45 sec")

# ── Slide 4: Beat 3 — The Result ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_label(s, "Beat 3 — The Result", Inches(0.9))
add_heading(s, "The Result", Inches(1.3), size=32)
add_body(s,
    '"Here\'s what we built."',
    Inches(2.3), size=17, italic=True)
add_note(s, "Let the video carry this. Say almost nothing. Point out one or two things if it feels natural.")
add_time_badge(s, "~90 sec")

# ── Slide 5: Beat 4 — Biggest Surprise ──────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_label(s, "Beat 4 — Biggest Surprise", Inches(0.9))
add_heading(s, "Biggest Surprise", Inches(1.3), size=32)
add_body(s,
    '"I also wanted it automatically installed on my phone every time I pushed code. '
    "I knew what that should look like — I build pipelines at work — but I had no idea what tools a game pipeline uses. "
    "I described the outcome. Claude came back with a complete solution: GitHub Actions, a Windows installer, "
    'a signed Android APK, Firebase pushing it to my phone. Thirty seconds after a git push."',
    Inches(2.3), size=16, italic=True)
add_note(s, "A screenshot here helps — the green Actions run or the Firebase notification on your phone.")
add_time_badge(s, "~60 sec")

# ── Slide 6: Beat 5 — The Closer ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_label(s, "Beat 5 — The Closer", Inches(0.9))
add_heading(s, "The Closer", Inches(1.3), size=32)
add_body(s,
    '"Three goals. None of them required me to already know how to do them."',
    Inches(2.5), size=20, italic=True)
add_note(s, "Stop there. That line does the work.")
add_time_badge(s, "~15 sec")

out = r"D:\Claude\Critter Quitters\claude\lightning_talk.pptx"
prs.save(out)
print(f"Saved: {out}")
